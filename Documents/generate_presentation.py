# pyright: reportUnknownParameterType=false, reportUnknownMemberType=false, reportUnknownVariableType=false, reportUnknownArgumentType=false, reportMissingParameterType=false, reportMissingTypeArgument=false
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
IMAGE_DIR = ROOT / "Image"
OUTPUT = ROOT / "Documents" / "DengAI_Project_Presentation.pptx"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)


def add_image_slide(prs, title, image_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.3), width=Inches(12.1))
    else:
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(2))
        tx.text_frame.text = f"Image not found: {image_path.name}"

    if caption:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11), Inches(0.6))
        box.text_frame.text = caption


def add_two_images_slide(prs, title, image1, image2, cap1="", cap2=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left_x, right_x = Inches(0.5), Inches(6.7)
    y, w = Inches(1.3), Inches(5.8)

    if image1.exists():
        slide.shapes.add_picture(str(image1), left_x, y, width=w)
    else:
        tx = slide.shapes.add_textbox(left_x, Inches(3.2), w, Inches(1))
        tx.text_frame.text = f"Missing: {image1.name}"

    if image2.exists():
        slide.shapes.add_picture(str(image2), right_x, y, width=w)
    else:
        tx = slide.shapes.add_textbox(right_x, Inches(3.2), w, Inches(1))
        tx.text_frame.text = f"Missing: {image2.name}"

    if cap1:
        c1 = slide.shapes.add_textbox(left_x, Inches(6.7), w, Inches(0.6))
        c1.text_frame.text = cap1
    if cap2:
        c2 = slide.shapes.add_textbox(right_x, Inches(6.7), w, Inches(0.6))
        c2.text_frame.text = cap2


def add_table_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows, cols = 6, 2
    table = slide.shapes.add_table(rows, cols, Inches(1.0), Inches(1.8), Inches(11.2), Inches(3.8)).table
    table.columns[0].width = Inches(5.2)
    table.columns[1].width = Inches(6.0)

    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"

    metrics = [
        ("Best Model", "Tuned XGBoost"),
        ("Validation MAE", "12.16"),
        ("Validation R²", "0.745"),
        ("Cross-Validation MAE (mean)", "10.71"),
        ("Training MAE", "3.59"),
    ]

    for r, (k, v) in enumerate(metrics, start=1):
        table.cell(r, 0).text = k
        table.cell(r, 1).text = v


def build_presentation():
    prs = Presentation()

    add_title_slide(
        prs,
        "DengAI: Predicting Disease Spread",
        "Comprehensive ML project presentation\nPrepared for GitHub portfolio",
    )

    add_bullets_slide(
        prs,
        "Problem Statement",
        [
            "Goal: Predict weekly dengue cases from climate/environment signals.",
            "Locations: San Juan (sj) and Iquitos (iq).",
            "Impact: Supports early warning and resource planning.",
        ],
    )

    add_bullets_slide(
        prs,
        "Dataset & Features",
        [
            "Training rows: 1456 | Test rows: 416",
            "Target: total_cases (weekly dengue cases)",
            "Feature groups: temperature, humidity, precipitation, NDVI, date-based fields",
            "Engineered features: month, day_of_year, city/season indicators",
        ],
    )

    add_bullets_slide(
        prs,
        "Methodology",
        [
            "Data cleaning and missing-value imputation",
            "Exploratory data analysis and city-level trend analysis",
            "Model comparison across 10 regression algorithms",
            "Hyperparameter tuning with GridSearchCV for XGBoost",
            "Validation + cross-validation performance checks",
        ],
    )

    add_image_slide(
        prs,
        "Target Variable and City Distribution",
        IMAGE_DIR / "01_target_analysis.png",
        "Cases are right-skewed, with higher spread and outliers in San Juan.",
    )

    add_image_slide(
        prs,
        "Model Comparison",
        IMAGE_DIR / "02_model_comparison_mae.png",
        "XGBoost shows the lowest validation MAE among tested models.",
    )

    add_image_slide(
        prs,
        "Feature Importance (Tuned XGBoost)",
        IMAGE_DIR / "03_feature_importance.png",
        "Environmental and seasonal variables drive prediction quality.",
    )

    add_two_images_slide(
        prs,
        "Model Diagnostics",
        IMAGE_DIR / "04_validation_actual_vs_predicted.png",
        IMAGE_DIR / "06_residual_diagnostics.png",
        "Actual vs predicted (validation)",
        "Residual distribution and residual-vs-predicted",
    )

    add_two_images_slide(
        prs,
        "Seasonality and Long-Term Trends",
        IMAGE_DIR / "05_city_month_heatmap.png",
        IMAGE_DIR / "07_yearly_trend_by_city.png",
        "Monthly average cases by city",
        "Yearly trend by city",
    )

    add_table_slide(prs, "Performance Summary")

    add_bullets_slide(
        prs,
        "Submission Output",
        [
            "Submission file generated: Data/submission.csv",
            "Rows: 416, matching required format",
            "Prediction stats: min=1, max=108, mean=22.76",
            "Ready for competition upload / benchmark scoring",
        ],
    )

    add_bullets_slide(
        prs,
        "Conclusion & Next Steps",
        [
            "Tuned XGBoost achieved strong validation performance.",
            "Model is suitable as a decision-support tool for dengue risk planning.",
            "Next: add lag features, city-specific models, and time-series CV.",
            "Deploy with weekly retraining and monitoring for concept drift.",
        ],
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Presentation created: {OUTPUT}")


if __name__ == "__main__":
    build_presentation()
